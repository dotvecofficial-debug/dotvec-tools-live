#!/usr/bin/env python3
import argparse
import csv
import difflib
import html
import os
import sys


def require(module, package_name=None):
    try:
        return __import__(module)
    except ImportError as exc:
        raise RuntimeError(f"Python package {package_name or module} is not installed") from exc


def redact(input_path, output_path, page_number, rect):
    fitz = require('fitz', 'PyMuPDF')
    doc = fitz.open(input_path)
    index = max(0, min(len(doc) - 1, page_number - 1))
    page = doc[index]
    x, y, width, height = rect
    area = fitz.Rect(x, y, x + width, y + height)
    page.add_redact_annot(area, fill=(0, 0, 0))
    page.apply_redactions()
    doc.save(output_path, garbage=4, deflate=True)


def compare(first, second, output_path):
    fitz = require('fitz', 'PyMuPDF')
    def extract(path):
        doc = fitz.open(path)
        return '\n\n'.join(f"--- Page {i+1} ---\n{page.get_text('text')}" for i, page in enumerate(doc))
    a = extract(first).splitlines(keepends=True)
    b = extract(second).splitlines(keepends=True)
    diff = difflib.unified_diff(a, b, fromfile=os.path.basename(first), tofile=os.path.basename(second), n=3)
    with open(output_path, 'w', encoding='utf-8') as handle:
        handle.writelines(diff)


def pdf_to_word(input_path, output_path):
    fitz = require('fitz', 'PyMuPDF')
    docx = require('docx', 'python-docx')
    pdf = fitz.open(input_path)
    document = docx.Document()
    for index, page in enumerate(pdf):
        if index:
            document.add_page_break()
        document.add_heading(f'Page {index + 1}', level=2)
        text = page.get_text('text').strip()
        if text:
            for paragraph in text.split('\n\n'):
                document.add_paragraph(paragraph.strip())
        else:
            document.add_paragraph('[No selectable text found on this page]')
    document.save(output_path)


def pdf_to_excel(input_path, output_path):
    fitz = require('fitz', 'PyMuPDF')
    openpyxl = require('openpyxl')
    pdf = fitz.open(input_path)
    workbook = openpyxl.Workbook()
    workbook.remove(workbook.active)
    for index, page in enumerate(pdf):
        sheet = workbook.create_sheet(title=f'Page {index + 1}'[:31])
        lines = [line.strip() for line in page.get_text('text').splitlines() if line.strip()]
        for row_index, line in enumerate(lines, start=1):
            cells = [part.strip() for part in line.split('\t')] if '\t' in line else [line]
            for column_index, value in enumerate(cells, start=1):
                sheet.cell(row=row_index, column=column_index, value=value)
        sheet.freeze_panes = 'A1'
    if not workbook.sheetnames:
        workbook.create_sheet('Page 1')
    workbook.save(output_path)


def pdf_to_powerpoint(input_path, output_path):
    fitz = require('fitz', 'PyMuPDF')
    pptx = require('pptx', 'python-pptx')
    from pptx.util import Inches
    pdf = fitz.open(input_path)
    presentation = pptx.Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    temp_files = []
    try:
        for index, page in enumerate(pdf):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            image_path = f"{output_path}.page-{index + 1}.png"
            pix.save(image_path)
            temp_files.append(image_path)
            slide = presentation.slides.add_slide(blank)
            slide.shapes.add_picture(image_path, 0, 0, width=presentation.slide_width, height=presentation.slide_height)
        if len(presentation.slides) == 0:
            presentation.slides.add_slide(blank)
        presentation.save(output_path)
    finally:
        for file_path in temp_files:
            try:
                os.remove(file_path)
            except OSError:
                pass


def flatten(input_path, output_path):
    fitz = require('fitz', 'PyMuPDF')
    doc = fitz.open(input_path)
    try:
        doc.bake(annots=True, widgets=True)
    except TypeError:
        doc.bake()
    doc.save(output_path, garbage=4, deflate=True)


def docx_to_text(input_path, output_path):
    docx = require('docx', 'python-docx')
    document = docx.Document(input_path)
    with open(output_path, 'w', encoding='utf-8') as handle:
        for paragraph in document.paragraphs:
            handle.write(paragraph.text + '\n')
        for table in document.tables:
            handle.write('\n')
            for row in table.rows:
                handle.write('\t'.join(cell.text for cell in row.cells) + '\n')


def xlsx_to_csv(input_path, output_path):
    openpyxl = require('openpyxl')
    workbook = openpyxl.load_workbook(input_path, data_only=True, read_only=True)
    sheet = workbook[workbook.sheetnames[0]]
    with open(output_path, 'w', newline='', encoding='utf-8-sig') as handle:
        writer = csv.writer(handle)
        for row in sheet.iter_rows(values_only=True):
            writer.writerow(['' if value is None else value for value in row])


def csv_to_xlsx(input_path, output_path):
    openpyxl = require('openpyxl')
    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = 'Data'
    with open(input_path, newline='', encoding='utf-8-sig') as handle:
        for row in csv.reader(handle):
            sheet.append(row)
    workbook.save(output_path)


def main():
    parser = argparse.ArgumentParser()
    sub = parser.add_subparsers(dest='command', required=True)

    p = sub.add_parser('redact')
    p.add_argument('input'); p.add_argument('output'); p.add_argument('--page', type=int, default=1); p.add_argument('--rect', nargs=4, type=float, required=True)
    p = sub.add_parser('compare'); p.add_argument('first'); p.add_argument('second'); p.add_argument('output')
    p = sub.add_parser('pdf-to-word'); p.add_argument('input'); p.add_argument('output')
    p = sub.add_parser('pdf-to-excel'); p.add_argument('input'); p.add_argument('output')
    p = sub.add_parser('pdf-to-powerpoint'); p.add_argument('input'); p.add_argument('output')
    p = sub.add_parser('flatten'); p.add_argument('input'); p.add_argument('output')
    p = sub.add_parser('docx-to-text'); p.add_argument('input'); p.add_argument('output')
    p = sub.add_parser('xlsx-to-csv'); p.add_argument('input'); p.add_argument('output')
    p = sub.add_parser('csv-to-xlsx'); p.add_argument('input'); p.add_argument('output')

    args = parser.parse_args()
    if args.command == 'redact': redact(args.input, args.output, args.page, args.rect)
    elif args.command == 'compare': compare(args.first, args.second, args.output)
    elif args.command == 'pdf-to-word': pdf_to_word(args.input, args.output)
    elif args.command == 'pdf-to-excel': pdf_to_excel(args.input, args.output)
    elif args.command == 'pdf-to-powerpoint': pdf_to_powerpoint(args.input, args.output)
    elif args.command == 'flatten': flatten(args.input, args.output)
    elif args.command == 'docx-to-text': docx_to_text(args.input, args.output)
    elif args.command == 'xlsx-to-csv': xlsx_to_csv(args.input, args.output)
    elif args.command == 'csv-to-xlsx': csv_to_xlsx(args.input, args.output)


if __name__ == '__main__':
    try:
        main()
    except Exception as exc:
        print(str(exc), file=sys.stderr)
        sys.exit(1)
